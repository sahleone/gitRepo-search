import os
import subprocess
from pptx import Presentation
import glob
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import glob
import os

import os
from PIL import Image
import numpy as np


from preprocess import normalize_text

def clone_repository(repo_url, target_dir):
    """
    Clones a Git repository into a specified target directory.

    Args:
        repo_url (str): URL of the Git repository to clone.
        target_dir (str): Path to the directory where the repository should be cloned.
    """
    try:
        # Create the target directory if it does not exist
        if not os.path.exists(target_dir):
            os.makedirs(target_dir)
        
        # Run the git clone command using subprocess.run
        result = subprocess.run(['git', 'clone', repo_url, target_dir], check=True, capture_output=True, text=True)
        
        # Print the result of the cloning process
        print(result.stdout)
        print(f"Repository cloned successfully into {target_dir}")
    
    except subprocess.CalledProcessError as e:
        # Print the error if the cloning process fails
        print(f"Failed to clone repository: {e.stderr}")
        
        
def extract_text_from_pptx_files(directory):
    # List to store tuples of (filename, text)
    presentations_text = []

    # Change the working directory to the specified directory
    os.chdir(directory)

    # Iterate over each pptx file in the directory
    for eachfile in glob.glob("*.pptx"):
        prs = Presentation(eachfile)
        slides_text = []

        # Iterate over each slide in the presentation
        for slide in prs.slides:
            slide_text_runs = []

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        slide_text_runs.append(run.text)
            
            # Normalize and join text for the slide
            normalized_slide_text = normalize_text(' '.join(slide_text_runs))
            slides_text.append(normalized_slide_text)
        
        # Join slide texts and append as a tuple to the list
        presentations_text.append((eachfile, '\n'.join(slides_text)))
    
    return presentations_text


def extract_text_from_pptx_files(directory):
    # List to store tuples of (filename, text)
    presentations_text = []

    # Change the working directory to the specified directory
    os.chdir(directory)

    # Iterate over each pptx file in the directory
    for eachfile in glob.glob("*.pptx"):
        prs = Presentation(eachfile)
        slides_text = []

        # Iterate over each slide in the presentation
        for slide in prs.slides:
            slide_text_runs = []

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        slide_text_runs.append(run.text)
            
            # Normalize and join text for the slide
            normalized_slide_text = normalize_text(' '.join(slide_text_runs))
            slides_text.append(normalized_slide_text)
        
        # Join slide texts and append as a tuple to the list
        presentations_text.append((eachfile, '\n'.join(slides_text)))
    
    return presentations_text

def __iter_picture_shapes(prs):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                yield shape

def extract_images_from_pptx_files(directory):
    presentations_images = []

    for eachfile in glob.glob(os.path.join(directory, "*.pptx")):
        prs = Presentation(eachfile)
        image_data_list = []

        for picture in __iter_picture_shapes(prs):
            image = picture.image
            image_bytes = image.blob
            image_data_list.append({
                'filename': eachfile,
                'image_data': image_bytes,
                'image_extension': image.ext
            })

        presentations_images.append((eachfile, image_data_list))
    
    return presentations_images




def load_images_from_folder(folder):
    allowed_extensions = {'.png', '.jpg', '.jpeg', '.gif', '.svg'}
    images = []
    filenames = []

    for filename in os.listdir(folder):
        ext = os.path.splitext(filename)[1].lower()
        if ext in allowed_extensions:
            img_path = os.path.join(folder, filename)
            img = Image.open(img_path)
            img_array = np.array(img)
            images.append(img_array)
            filenames.append(filename)
    
    return images, filenames
